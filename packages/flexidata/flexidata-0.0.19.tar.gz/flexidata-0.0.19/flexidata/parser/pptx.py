from flexidata.reader.handlers import FileHandler
from typing import Optional, List, BinaryIO, cast, Dict
from flexidata.utils.constants import FileReaderSource, FileType
from flexidata.text_processing.classification import TextType
from pptx import Presentation
from flexidata.text_processing.text_block import TextBlock, TextBlockMetadata
from flexidata.utils.decorators import validate_file_type_method
from pptx.shapes.graphfrm import GraphicFrame
from pptx.shapes.picture import Picture
from flexidata.utils.tables import (
    extract_inner_tables,
    remove_inner_tables,
    table_dict_to_html_table,
    table_dict_to_plain_text,
)
import numpy as np
import io
from flexidata.ocr.agent import OCRAgentFactory
from PIL import Image

class PptxParser(FileHandler):

    def __init__(
        self,
        file_path: Optional[str] = None,
        file_url: Optional[str] = None,
        source: FileReaderSource = FileReaderSource.LOCAL,
        extract_table: bool = False,
        extract_image: bool = False,
        bucket_name: Optional[str] = None,
        file_key: Optional[str] = None,
        table_dict_in_metadata: bool = False,
        separate_inner_tables: bool = True,
        **kwargs,
    ) -> None:
        super().__init__(source, file_path, file_url, bucket_name, file_key)
        self.extract_table = extract_table
        self.extract_image = extract_image
        self.table_dict_in_metadata = table_dict_in_metadata
        self.separate_inner_tables = separate_inner_tables
        self.text_extractable = False
        self.presentation = None
        self.page_number = 1
        self.text_blocks = []

    @validate_file_type_method([FileType.PPTX])
    def parse(self, extract_image: bool = False) -> List[TextBlock]:
        if extract_image:
            self.extract_image = extract_image
        self.presentation = Presentation(cast(BinaryIO, self.get_file_content()))

        for slide_index, slide in enumerate(self.presentation.slides):
            self.page_number = slide_index + 1
            self._process_shapes(slide.shapes)

        return self.text_blocks

    def _process_shapes(self, shapes):
        for shape in shapes:
            if shape.has_table and isinstance(shape, GraphicFrame):
                self._process_table(shape)
            elif shape.has_chart:
                self._process_chart(shape)
            elif shape.has_text_frame:
                self._process_text(shape)
            elif self.extract_image and isinstance(shape, Picture):
                self._process_image(shape)

    def _process_image(self, shape, format='list'):
        image_texts = []
        image_text = ''
        image_blob = shape.image.blob
        image = np.array(Image.open(io.BytesIO(image_blob)))
        ocr_factory = OCRAgentFactory()
        ocr_agent = ocr_factory.get_ocr_agent()
        elements = ocr_agent.image_to_text(
            np.array(image),
            lang="eng",
            page_number=self.page_number,
            file_path=self.get_file_path_by_source(),
            filetype=FileType.DOCX,
        )
        if format == 'plain':
            for element in elements:
                image_text += element.text + "\n"
        if format == 'list':
            image_texts.extend(elements)
            self.text_blocks.extend(image_texts)
        return image_text if format == 'plain' else image_texts

    def _process_table(self, shape):
        table_dict = self._convert_table_to_dict(shape.table)
        if self.separate_inner_tables:
            main_table = remove_inner_tables(table_dict)
            self._process_table_dict(main_table)
            inner_tables = extract_inner_tables(table_dict)
            for table in inner_tables:
                self._process_table_dict(table)
        else:
            self._process_table_dict(table_dict)

    def _process_table_dict(self, table_dict: Dict):
        """
        Processes a table dictionary to generate plain text and HTML representations,
        create metadata, and append a text block to the document's text blocks collection.

        Args:
            table_dict (Dict): The dictionary representation of a table that includes data
                               necessary to render text and HTML formats.

        Description:
            The function converts a table dictionary into plain text and HTML formats,
            generates metadata based on the HTML and potentially the table dictionary itself,
            and then creates and appends a new text block to the document's collection of text blocks.
        """
        text_table = table_dict_to_plain_text(table_dict)
        html_table = table_dict_to_html_table(table_dict)

        metadata = self._create_metadata(
            html_table, table_dict if self.table_dict_in_metadata else None
        )

        text_block = TextBlock(text_table, metadata, TextType.TABLES)
        self.text_blocks.append(text_block)

    def _convert_table_to_dict(
        self, table, table_dict: Optional[Dict[str, List]] = None
    ):
        if table_dict is None:
            table_dict = {"rows": []}
            inner_tables = []

        for row in table.rows:
            columns = {}
            for column_index, cell in enumerate(row.cells):
                column = {'value': '', 'row_span': 0, 'col_span': 0}
                column['value'] = cell.text
                columns[f"column_{column_index}"] = column

            table_dict["rows"].append(columns)
        return table_dict

    def _process_chart(self, shape):
        pass

    def _process_text(self, shape):
        for paragraph in shape.text_frame.paragraphs:
            text_type = None
            text = paragraph.text
            if text.strip() == "":
                continue

            if self._is_bulleted_paragraph(paragraph):
                text_type = TextType.LIST_ITEM
            metadata = self._create_metadata("")
            text_block = TextBlock(text, metadata, text_type)
            self.text_blocks.append(text_block)

    def _create_metadata(self, html_table: str, table_dict: Optional[Dict] = None):

        metadata_details = {
            "text_as_html": html_table,
            "file_path": self.get_file_path_by_source(),
            "languages": ["eng"],
            "filetype": FileType.PPTX,
            "page_number": self.page_number,
            "extraction_source": "pptx",
        }
        if table_dict is not None:
            metadata_details["table_dict"] = table_dict

        return TextBlockMetadata(**metadata_details)

    def _is_bulleted_paragraph(self, paragraph) -> bool:
        return bool(paragraph._p.xpath("./a:pPr/a:buChar"))
